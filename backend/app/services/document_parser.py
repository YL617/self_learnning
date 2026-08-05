from __future__ import annotations

from pathlib import Path


def extract_text(file_path: str, file_type: str) -> str:
    suffix = file_type.lower()
    if suffix in {"pdf"}:
        from pypdf import PdfReader

        reader = PdfReader(file_path)
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    if suffix in {"docx"}:
        from docx import Document

        doc = Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    if suffix in {"pptx"}:
        from pptx import Presentation

        prs = Presentation(file_path)
        lines: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.append(shape.text.strip())
        return "\n".join(lines)
    if suffix in {"txt", "md"}:
        return Path(file_path).read_text(encoding="utf-8", errors="ignore")
    if suffix in {"png", "jpg", "jpeg"}:
        # OCR 引擎（如 Tesseract / 云 OCR）在第二阶段接入
        raise ValueError("图片 OCR 尚未启用，请先使用文本类文档")
    raise ValueError(f"不支持的文档类型：{suffix}")


def chunk_text(text: str, chunk_size: int = 800, overlap: int = 100) -> list[str]:
    text = text.strip()
    if not text:
        return []
    paragraphs = [p.strip() for p in text.splitlines() if p.strip()]
    chunks: list[str] = []
    buffer = ""
    for paragraph in paragraphs:
        if len(buffer) + len(paragraph) + 1 <= chunk_size:
            buffer = f"{buffer}\n{paragraph}" if buffer else paragraph
            continue
        if buffer:
            chunks.append(buffer)
        if len(paragraph) > chunk_size:
            start = 0
            while start < len(paragraph):
                end = min(start + chunk_size, len(paragraph))
                chunk = paragraph[start:end]
                chunks.append(chunk)
                start = end - overlap if end - overlap > start else end
        else:
            buffer = paragraph
    if buffer:
        chunks.append(buffer)
    return chunks
